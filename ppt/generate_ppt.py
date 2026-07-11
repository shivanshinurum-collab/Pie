import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def main():
    # Initialize
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors
    PRIMARY_COLOR = RGBColor(22, 38, 57)      # Slate Navy
    SECONDARY_COLOR = RGBColor(34, 139, 139)  # Teal Blue
    ACCENT_COLOR = RGBColor(192, 57, 43)      # Crimson Red
    TEXT_COLOR = RGBColor(60, 60, 60)         # Dark Gray
    BG_COLOR = RGBColor(248, 249, 250)        # Off-white background
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(230, 235, 240)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_header(slide, category, title):
        # Category Tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_left = cat_tf.margin_right = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = 'Arial'
        cat_p.font.size = Pt(14)  # Massively increased for large halls
        cat_p.font.bold = True
        cat_p.font.color.rgb = SECONDARY_COLOR
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_left = title_tf.margin_right = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.name = 'Georgia'
        title_p.font.size = Pt(40)  # Massively increased for large halls
        title_p.font.bold = True
        title_p.font.color.rgb = PRIMARY_COLOR

    def add_title_slide(prs):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, BG_COLOR)
        
        # Decorative left accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = SECONDARY_COLOR
        accent_bar.line.fill.background()
        
        # Title text box (left column)
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(6.0), Inches(3.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p1 = tf.paragraphs[0]
        p1.text = "METABOLIC BONE DISORDERS"
        p1.font.name = 'Georgia'
        p1.font.size = Pt(56)  # Huge title
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_COLOR
        p1.space_after = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = "A Clinical Guide to Osteoporosis, Osteomalacia, and Paget's Disease"
        p2.font.name = 'Arial'
        p2.font.size = Pt(28)  # Huge subtitle
        p2.font.color.rgb = SECONDARY_COLOR
        p2.space_after = Pt(45)
        
        p3 = tf.add_paragraph()
        p3.text = "Prepared by Clinical Endocrinology & Nursing Education Experts"
        p3.font.name = 'Arial'
        p3.font.size = Pt(18)  # Huge author info
        p3.font.italic = True
        p3.font.color.rgb = TEXT_COLOR

        # Title image on right
        image_path = "images/metabolic_bone_intro.png"
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(7.5), Inches(1.3), Inches(5.0), Inches(5.0))

    def add_image_slide(prs, category, title, bullets, image_path):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, BG_COLOR)
        add_slide_header(slide, category, title)
        
        # Left Text Column
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Line-wrap estimation to maximize font sizes dynamically
        total_lines = 0
        for bullet in bullets:
            text = bullet[0]
            level = bullet[1]
            # Width is 6 inches. At Pt(24), average characters per line is about 40
            char_per_line = 36 if level == 0 else 42
            lines = max(1, len(text) // char_per_line + 1)
            total_lines += lines

        # Select the absolute largest sizes that fit the text box safely
        if total_lines > 12:
            sz_0 = Pt(20)
            sz_1 = Pt(16)
            space = Pt(4)
        elif total_lines > 9:
            sz_0 = Pt(24)
            sz_1 = Pt(20)
            space = Pt(6)
        elif total_lines > 6:
            sz_0 = Pt(28)
            sz_1 = Pt(23)
            space = Pt(8)
        else:
            sz_0 = Pt(34)  # Giant fonts for low-density slides!
            sz_1 = Pt(28)
            space = Pt(12)
            
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet[0]
            p.font.name = 'Arial'
            p.font.size = sz_0 if bullet[1] == 0 else sz_1
            p.font.color.rgb = TEXT_COLOR
            p.level = bullet[1]
            p.space_after = space
            
            if len(bullet) > 2:
                if bullet[2] == 'accent':
                    p.font.color.rgb = ACCENT_COLOR
                    p.font.bold = True
                elif bullet[2] == 'primary':
                    p.font.color.rgb = PRIMARY_COLOR
                    p.font.bold = True

        # Right Image Column
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(7.1), Inches(1.8), Inches(5.4), Inches(4.9))
        else:
            # Fallback Box
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(1.8), Inches(5.4), Inches(4.9)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = LIGHT_GRAY
            placeholder.line.color.rgb = SECONDARY_COLOR
            tf_ph = placeholder.text_frame
            tf_ph.word_wrap = True
            p = tf_ph.paragraphs[0]
            p.text = f"\n\n\n\n[Illustration: {os.path.basename(image_path)}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = PRIMARY_COLOR
            p.font.bold = True

    def add_bullet_slide(prs, category, title, bullets):
        add_image_slide(prs, category, title, bullets, "images/metabolic_bone_intro.png")

    def add_table_slide(prs, category, title, headers, rows, image_path):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, BG_COLOR)
        add_slide_header(slide, category, title)
        
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(8.9)  # Leave room for image on right
        height = Inches(4.8)
        
        num_rows = len(rows) + 1
        num_cols = len(headers)
        
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table
        
        # Custom Column Widths
        table.columns[0].width = Inches(1.8) # Disease Name
        table.columns[1].width = Inches(2.6) # Primary Pathology
        table.columns[2].width = Inches(1.3) # Serum Ca & PO4
        table.columns[3].width = Inches(1.2) # Serum ALP
        table.columns[4].width = Inches(2.0) # Main Treatments
        
        # Format Headers (Huge text)
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            cell.text = header_text
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(16)  # Increased for massive visibility
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
        # Format Data Rows (Huge text)
        for row_idx, row_data in enumerate(rows):
            for col_idx, val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.fill.solid()
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                else:
                    cell.fill.fore_color.rgb = WHITE
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.font.name = 'Arial'
                p.font.size = Pt(13)  # Increased for massive visibility
                p.font.color.rgb = TEXT_COLOR
                # Alignments
                if col_idx == 0:
                    p.alignment = PP_ALIGN.LEFT
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.CENTER

        # Table Slide Right-side Image
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(9.9), Inches(1.8), Inches(2.8), Inches(4.8))
        else:
            # Fallback Box
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(9.9), Inches(1.8), Inches(2.8), Inches(4.8)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = LIGHT_GRAY
            placeholder.line.color.rgb = SECONDARY_COLOR
            tf_ph = placeholder.text_frame
            tf_ph.word_wrap = True
            p = tf_ph.paragraphs[0]
            p.text = f"\n\n\n\n[Illustration: {os.path.basename(image_path)}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = PRIMARY_COLOR
            p.font.bold = True

    def add_nursing_care_plan_slide(prs, title, assessment, diagnosis, goal, interventions, evaluation):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, BG_COLOR)
        add_slide_header(slide, "Nursing Care Plan", title)
        
        # Draw 3 cards (rectangles) for visual structure
        # Card 1: Diagnosis & Assessment
        card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.9))
        card1.fill.solid()
        card1.fill.fore_color.rgb = WHITE
        card1.line.color.rgb = LIGHT_GRAY
        tf1 = card1.text_frame
        tf1.word_wrap = True
        tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = Inches(0.15)
        
        p = tf1.paragraphs[0]
        p.text = "ASSESSMENT"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(6)
        
        for item in assessment:
            p = tf1.add_paragraph()
            p.text = item
            p.font.name = 'Arial'
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(3)
            
        p = tf1.add_paragraph()
        p.text = "\nNURSING DIAGNOSIS"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(6)
        
        p = tf1.add_paragraph()
        p.text = diagnosis
        p.font.name = 'Arial'
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Card 2: Goals & Evaluation
        card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(1.8), Inches(3.7), Inches(4.9))
        card2.fill.solid()
        card2.fill.fore_color.rgb = WHITE
        card2.line.color.rgb = LIGHT_GRAY
        tf2 = card2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = Inches(0.15)
        
        p = tf2.paragraphs[0]
        p.text = "PLANNING & GOALS"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(6)
        
        for item in goal:
            p = tf2.add_paragraph()
            p.text = item
            p.font.name = 'Arial'
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(3)
            
        p = tf2.add_paragraph()
        p.text = "\nEVALUATION"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(6)
        
        for item in evaluation:
            p = tf2.add_paragraph()
            p.text = item
            p.font.name = 'Arial'
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(3)

        # Card 3: Interventions
        card3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(1.8), Inches(3.9), Inches(4.9))
        card3.fill.solid()
        card3.fill.fore_color.rgb = WHITE
        card3.line.color.rgb = LIGHT_GRAY
        tf3 = card3.text_frame
        tf3.word_wrap = True
        tf3.margin_left = tf3.margin_right = tf3.margin_top = tf3.margin_bottom = Inches(0.15)
        
        p = tf3.paragraphs[0]
        p.text = "NURSING INTERVENTIONS"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(6)
        
        for item in interventions:
            p = tf3.add_paragraph()
            p.text = item
            p.font.name = 'Arial'
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(3)

    # --- BUILD PRESENTATION ---

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: Intro Definition (verbatim from PDF Page 1)
    add_image_slide(prs, "General Overview", "Introduction & Definition of Metabolic Bone Disorders", [
        ["Metabolic bone disorders are a group of conditions that affect the strength, structure, and function of bones due to abnormalities in mineral metabolism, especially calcium, phosphorus, and vitamin D.", 0],
        ["These disorders typically result from disruptions in the processes that regulate bone formation (osteoblast activity) and bone resorption (osteoclast activity).", 0],
        ["Bones are dynamic tissues that constantly undergo remodeling to maintain their integrity.", 0],
        ["- This process depends on adequate nutrition, hormonal balance, and proper functioning of organs such as the kidneys, intestines, and endocrine glands.", 1],
        ["- When these systems are disturbed, it can lead to weakened bones, deformities, pain, and increased risk of fractures.", 1],
        ["Common examples of metabolic bone disorders include:", 0, "primary"],
        ["- Osteoporosis – characterized by decreased bone mass and fragility", 1],
        ["- Osteomalacia – softening of bones due to vitamin D deficiency", 1],
        ["- Rickets – defective bone mineralization in children", 1],
        ["- Paget's disease of bone – abnormal bone destruction and regrowth", 1]
    ], "images/metabolic_bone_overview.png")

    # Slide 3: Anatomy: Structure of Bone (verbatim from PDF Page 1 & 2)
    add_image_slide(
        prs,
        "Anatomy of Bone",
        "Structure of Bone",
        [
            ["A typical long bone consists of:", 0, "primary"],
            ["- Diaphysis: Shaft of the bone made mainly of compact bone.", 1],
            ["- Epiphysis: Ends of the bone containing spongy (cancellous) bone.", 1],
            ["- Metaphysis: Region between the diaphysis and epiphysis; contains the growth plate in children.", 1],
            ["- Periosteum: Outer fibrous membrane rich in blood vessels and nerves.", 1],
            ["- Endosteum: Inner membrane lining the medullary cavity.", 1],
            ["- Medullary cavity: Contains yellow bone marrow in adults.", 1],
            ["- Articular cartilage: Covers the ends of bones and reduces friction at joints.", 1]
        ],
        "images/bone_remodeling.png"
    )

    # Slide 4: Anatomy: Types of Bone (verbatim from PDF Page 2)
    add_image_slide(
        prs,
        "Anatomy of Bone",
        "Types of Bone",
        [
            ["Bones are classified into five distinct types based on shape and location:", 0, "primary"],
            ["- Long bones (Femur, Humerus)", 1],
            ["- Short bones (Carpal bones)", 1],
            ["- Flat bones (Skull, Sternum)", 1],
            ["- Irregular bones (Vertebrae)", 1],
            ["- Sesamoid bones (Patella)", 1]
        ],
        "images/bone_types.png"
    )

    # Slide 5: Anatomy: Bone Cells & Matrix (verbatim from PDF Page 2)
    add_image_slide(
        prs,
        "Anatomy of Bone",
        "Bone Cells & Matrix",
        [
            ["3. Bone Cells:", 0, "primary"],
            ["- Osteoblasts: Bone-forming cells that synthesize bone matrix.", 1],
            ["- Osteocytes: Mature bone cells that maintain bone tissue.", 1],
            ["- Osteoclasts: Large multinucleated cells responsible for bone resorption.", 1],
            ["4. Bone Matrix:", 0, "accent"],
            ["- Organic Matrix (≈35%): Type I collagen, Proteoglycans, Glycoproteins.", 1],
            ["- Inorganic Matrix (≈65%): Calcium, Phosphate, Hydroxyapatite crystals, Magnesium, Fluoride.", 1]
        ],
        "images/bone_matrix_structure.png"
    )

    # Slide 6: Physiology: Functions & Remodeling Cycle (verbatim from PDF Page 2)
    add_image_slide(
        prs,
        "Physiology of Bone",
        "Functions & Remodeling Cycle",
        [
            ["Functions of Bone:", 0, "primary"],
            ["- Provides support and body shape.", 1],
            ["- Protects vital organs.", 1],
            ["- Acts as a lever for movement.", 1],
            ["- Stores calcium and phosphate.", 1],
            ["- Produces blood cells (hematopoiesis).", 1],
            ["- Stores fat in yellow marrow.", 1],
            ["Bone Remodeling (continuous process consisting of):", 0, "accent"],
            ["- Activation – Osteoclasts are stimulated.", 1],
            ["- Resorption – Old bone is broken down.", 1],
            ["- Reversal – Bone surface is prepared.", 1],
            ["- Formation – Osteoblasts form new bone.", 1],
            ["- Mineralization – Calcium and phosphate are deposited.", 1]
        ],
        "images/bone_remodeling.png"
    )

    # Slide 7: Physiology: Homeostasis & Hormonal Control (verbatim from PDF Page 2 & 3)
    add_image_slide(
        prs,
        "Physiology of Bone",
        "Calcium and Phosphate Homeostasis",
        [
            ["Normal Serum Mineral Levels:", 0, "primary"],
            ["- Normal Serum Calcium: 8.5–10.5 mg/dL", 1],
            ["- Normal Serum Phosphate: 2.5–4.5 mg/dL", 1],
            ["Hormonal Regulation:", 0, "accent"],
            ["1. Parathyroid Hormone (PTH): Increases blood calcium. Increases bone resorption, increases renal calcium reabsorption, decreases phosphate reabsorption, activates vitamin D.", 1],
            ["2. Vitamin D (Calcitriol): Increases intestinal absorption of calcium and phosphate; promotes bone mineralization.", 1],
            ["3. Calcitonin: Inhibits osteoclast activity, reduces bone resorption, and lowers blood calcium.", 1]
        ],
        "images/vit_d_metabolism.png"
    )

    # Slide 8: Osteoporosis: Definition (verbatim from PDF Page 3)
    add_image_slide(
        prs,
        "Osteoporosis",
        "Introduction & Definition",
        [
            ["Clinical Context:", 0, "primary"],
            ["- Most common metabolic bone disease.", 1],
            ["- One of the most prevalent conditions associated with aging.", 1],
            ["Definition:", 0, "accent"],
            ["- Osteoporosis is a disorder characterized by low bone mass and bone tissue degeneration, which increases the risk of fractures.", 1]
        ],
        "images/osteoporosis_bone.png"
    )

    # Slide 9: Osteoporosis: Etiology & Risk Factors (verbatim from PDF Page 3)
    add_image_slide(
        prs,
        "Osteoporosis",
        "Etiology & Risk Factors",
        [
            ["1. Age is a risk factor for osteoporosis, especially in postmeno-pausal women.", 0],
            ["2. Gender: Women are more likely than males to acquire osteoporosis.", 0],
            ["3. Hormonal factors: Both low testosterone levels in men and low estrogen levels during menopause contribute to bone loss.", 0],
            ["4. The risk of osteoporosis is increased by a family history of the disease.", 0],
            ["5. Lack of physical activity, smoking, excessive alcohol con-sumption, and inadequate calcium and vitamin D intake are all risk factors for osteoporosis.", 0],
            ["6. Medical diseases and drugs (eg, RA and corticosteroids) can raise the risk of osteoporosis.", 0, "accent"]
        ],
        "images/osteoporosis_spine.png"
    )

    # Slide 10: Osteoporosis: Pathophysiology (verbatim from PDF Page 3 & 4)
    add_image_slide(
        prs,
        "Osteoporosis",
        "Pathophysiology",
        [
            ["1. Bone remodeling: Bones are constantly remodeling, in which old bone is broken down and replaced by new bone tissue.", 0],
            ["2. Remodeling imbalance: There is an imbalance between bone resorption (breakdown) and bone production in osteoporosis, resulting in a net loss of bone mass.", 0],
            ["3. Increased bone resorption: Osteoclast activity is increased, resulting in excessive bone resorption.", 0, "accent"],
            ["4. Reduced bone production: Osteoblasts, which are responsible for bone production, are unable to successfully manufacture new bone tissue.", 0],
            ["5. As a result of the microarchitectural changes, the bone becomes porous, fragile, and prone to fractures.", 0]
        ],
        "images/bone_remodeling.png"
    )

    # Slide 11: Osteoporosis: Clinical Manifestations (verbatim from PDF Page 4)
    add_image_slide(
        prs,
        "Osteoporosis",
        "Clinical Manifestations",
        [
            ["1. Osteoporosis is frequently asymptomatic until a fracture occurs. Fractures are most typically found in the spine, hip, and wrist, although they can also occur in other bones.", 0, "accent"],
            ["2. Gradual decrease of height caused by spinal compression fractures.", 0],
            ["3. Chronic back pain: Spinal fractures can cause chronic back pain.", 0],
            ["4. Vertebral fractures can result in a forward curve of the spine (kyphosis) and a stooped posture.", 0, "accent"]
        ],
        "images/osteoporosis_spine.png"
    )

    # Slide 12: Osteoporosis: Diagnostic Tests (verbatim from PDF Page 4)
    add_image_slide(
        prs,
        "Osteoporosis",
        "Diagnostic Tests",
        [
            ["1. DXA analyzes bone mineral density (BMD) and helps diag-Pro osteoporosis and assess fracture risk.", 0, "accent"],
            ["2. Fracture risk assessment: Tools such as FRAX quantify the likelihood of a major osteoporotic fracture during a given time period based on a variety of risk variables.", 0],
            ["3. Blood tests can analyze calcium and vitamin D levels, as well as underlying medical disorders that contribute to osteoporosis.", 0]
        ],
        "images/osteoporosis_bone.png"
    )

    # Slide 13: Osteoporosis: Medical & Surgical Management (verbatim from PDF Page 4 & 5)
    add_image_slide(
        prs,
        "Osteoporosis",
        "Medical & Surgical Management",
        [
            ["Medical Management:", 0, "primary"],
            ["- 1. Encourage regular weight-bearing exercise, a balanced diet rich in calcium and vitamin D, and abstinence from smoking and excessive alcohol consumption.", 1],
            ["- 2. To prevent bone loss, raise bone density, and minimize fracture risk; medications such as bisphosphonates, selective estrogen receptor modulators (SERMs), denosumab, teri-paratide, or abaloparatide may be administered.", 1],
            ["- 3. Fall prevention: Implement fall prevention measures such as home safety modifications and balance exercises.", 1],
            ["Surgical Management:", 0, "accent"],
            ["- 1. Vertebroplasty and kyphoplasty are procedures that use cement or inflatable balloons to fix spine fractures.", 1],
            ["- 2. Joint replacement surgery: Joint replacement surgery may be considered if osteoporosis-related fractures severely impair the hip or knee joints.", 1]
        ],
        "images/metabolic_bone_overview.png"
    )

    # Slide 14: Osteoporosis: Nursing - Education & Nutrition (verbatim from PDF Page 5)
    add_image_slide(
        prs,
        "Osteoporosis: Nursing",
        "Education & Nutritional Support",
        [
            ["Education:", 0, "primary"],
            ["- 1. Provide comprehensive education about osteoporosis, including its causes, risk factors, prevention strategies, and management.", 1],
            ["- 2. Teach patients the benefits of a calcium- and vitamin D-rich diet for bone health.", 1],
            ["- 3. Patients should be educated on weight-bearing activities, bal-ance exercises, and fall prevention measures.", 1],
            ["- 4. Stress the need of drug adherence as well as the potential side effects of osteoporosis treatments.", 1, "accent"],
            ["- 5. Provide resources for osteoporosis support groups, educa-tional materials, and community events.", 1],
            ["Nutritional Support:", 0, "accent"],
            ["- 1. Work with dietitians to create a specific diet plan that focuses on calcium- and vitamin D-rich foods.", 1],
            ["- 2. Educate patients on calcium-rich foods such as dairy prod-ucts, leafy green vegetables, and fortified meals.", 1],
            ["- 3. Discuss how vitamin D and sunshine exposure affect calcium absorption.", 1],
            ["- 4. Assess patients' nutritional condition and, if necessary, sug-gest for supplementation.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 15: Osteoporosis: Nursing - Exercise Promotion (verbatim from PDF Page 5)
    add_image_slide(
        prs,
        "Osteoporosis: Nursing",
        "Exercise Promotion",
        [
            ["Exercise:", 0, "primary"],
            ["- 1. Encourage patients to participate in weight-bearing exercises such as walking, dancing, or stair climbing to enhance bone strength.", 1],
            ["- 2. Give instructions on suitable exercise practices, as well as the frequency and duration of exercise sessions.", 1],
            ["- 3. Collaborate with physical therapists to create tailored exer-cise programs depending on the talents and limitations of the patients.", 1],
            ["- 4. Stress the need for regular physical activity in maintaining bone health and preventing falls.", 1, "accent"]
        ],
        "images/osteoporosis_spine.png"
    )

    # Slide 16: Osteoporosis: Nursing - Fall & Med Management (verbatim from PDF Page 5 & 6)
    add_image_slide(
        prs,
        "Osteoporosis: Nursing",
        "Fall & Medication Management",
        [
            ["Fall Prevention:", 0, "primary"],
            ["- 1. Evaluate the patient's living settings for potential risks and provide recommendations for changes to ensure his or her safety.", 1],
            ["- 2. Patients should be educated on fall prevention techniques such as reducing tripping hazards, using assistive equipment, and wearing suitable footwear.", 1],
            ["- 3. Encourage the installation of grab bars and the use of nonslip mats or rugs in bathrooms.", 1],
            ["- 4. Teach patients skills for increasing balance and lowering the chance of falling, such as tai chi or excellent posture.", 1],
            ["Medication Management:", 0, "accent"],
            ["- 1. Educate patients on the purpose, dosage, administration, and potential side effects of their osteoporosis drugs.", 1],
            ["- 2. Insist on drug adherence and timely follow-up with health care providers.", 1],
            ["- 3. Inform patients about potential drug interactions and the sig-nificance of alerting their health care professionals about all medications they are taking.", 1],
            ["- 4. Patients should be monitored for pharmaceutical side effects and their response to treatment, and any concerns should be reported to the health care team.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 17: Osteoporosis: Nursing - Counseling & Support (verbatim from PDF Page 6)
    add_image_slide(
        prs,
        "Osteoporosis: Nursing",
        "Counseling & Emotional Support",
        [
            ["Counseling and Emotional Support:", 0, "primary"],
            ["- 1. Provide emotional support and counseling to patients to address their worries, concerns, and emotional impact of osteoporosis.", 1],
            ["- 2. Encourage patients to talk about their feelings and aid them in developing coping techniques.", 1, "accent"],
            ["- 3. Facilitate access to further psychological assistance through support groups, counseling services, or community resources.", 1]
        ],
        "images/osteoporosis_bone.png"
    )

    # Slide 18: Osteomalacia: Definition & Etiology (verbatim from PDF Page 6 & 7)
    add_image_slide(
        prs,
        "Osteomalacia",
        "Definition & Etiology",
        [
            ["Definition:", 0, "primary"],
            ["- Osteomalacia is a bone softening and weakening condition caused by poor mineralization of the bone matrix. It is primarily caused by a lack of vitamin D or malabsorption.", 1, "accent"],
            ["Etiology:", 0, "primary"],
            ["- 1. Vitamin D deficiency: Inadequate intake of vitamin D through diet or limited exposure to sunshine, both of which are required for vitamin D production in the skin.", 1],
            ["- 2. Conditions that impede vitamin D and calcium absorption. such as celiac disease, inflammatory bowel disease, or gastric bypass surgery, are examples of malabsorption.", 1],
            ["- 3. Medications: Certain drugs, such as anticonvulsants and ght cocorticoids (glucocorticoids), can interfere with vitamin D metabolism over time.", 1],
            ["- 4. Renal dysfunction: Chronic renal illness can decrease vitamin D activation and calcium absorption.", 1]
        ],
        "images/osteomalacia_legs.png"
    )

    # Slide 19: Osteomalacia: Pathophysiology (verbatim from PDF Page 7)
    add_image_slide(
        prs,
        "Osteomalacia",
        "Pathophysiology",
        [
            ["1. Vitamin D metabolism: Vitamin D is converted in the kidneys to its active form (calcitriol), which is required for calcium absorption from the intestines.", 0],
            ["2. Calcium homeostasis: Adequate quantities of vitamin D and calcium are required for the body to sustain calcium homeostasis.", 0],
            ["3. Impaired mineralization: Vitamin D insufficiency or malabsorption reduces dietary calcium absorption, resulting in low serum calcium levels.", 0, "accent"],
            ["4. Secondary hyperparathyroidism occurs when low calcium levels encourage the release of parathyroid hormone (PTH), resulting in enhanced bone resorption and calcium mobilization from bones.", 0],
            ["5. Inadequate mineralization of the bone matrix results in the accumulation of unmineralized osteoid, resulting in weaker and soft bones.", 0, "accent"]
        ],
        "images/bone_matrix_structure.png"
    )

    # Slide 20: Osteomalacia: Clinical Manifestations (verbatim from PDF Page 7 & 8)
    add_image_slide(
        prs,
        "Osteomalacia",
        "Clinical Manifestations",
        [
            ["1. Dull, agonizing bone pain, especially in the hips, lower back. and legs.", 0],
            ["2. Muscle weakness: Fatigue and generalized muscle weakness.", 0],
            ["3. Skeletal deformities include leg bowing, kyphosis (excessive curvature of the upper spine), and short stature.", 0, "accent"],
            ["4. Fracture vulnerability: Increased vulnerability to fractures, particularly in weight-bearing bones.", 0],
            ["5. Walking difficulties: Gait irregularities and difficulties doing daily activities as a result of discomfort and muscle weakness.", 0, "accent"]
        ],
        "images/osteomalacia_legs.png"
    )

    # Slide 21: Osteomalacia: Diagnostics & Imaging (verbatim from PDF Page 8)
    add_image_slide(
        prs,
        "Osteomalacia",
        "Diagnostic Tests & Imaging",
        [
            ["Diagnostic Tests:", 0, "primary"],
            ["- 1. Serum calcium and phosphorus levels: Low serum calcium and phosphate levels.", 1],
            ["- 2. Elevated levels of ALP due to accelerated bone turnover.", 1, "accent"],
            ["- 3. 25-Hydroxyvitamin D: Low 25-hydroxyvitamin D levels indicate vitamin D insufficiency.", 1],
            ["- 4. PTH: Elevated PTH levels as a result of secondary hyperparathyroidism.", 1],
            ["Imaging:", 0, "accent"],
            ["- 1. X-rays may reveal pseudo fractures or Looser zones (radiolu-cent lines) in the bones.", 1],
            ["- 2. DXA: Assesses the severity of osteomalacia by measuring BMD.", 1]
        ],
        "images/vit_d_metabolism.png"
    )

    # Slide 22: Osteomalacia: Medical & Surgical Management (verbatim from PDF Page 8)
    add_image_slide(
        prs,
        "Osteomalacia",
        "Medical & Surgical Management",
        [
            ["Medical Management:", 0, "primary"],
            ["- 1. Oral vitamin D supplementation is used to treat vitamin D insufficiency or malabsorption.", 1, "accent"],
            ["- 2. Oral calcium supplements are used to guarantee appropriate calcium intake.", 1],
            ["- 3. Treatment of underlying disorders that contribute to vitamin D deficiency or malabsorption.", 1],
            ["- 4. Changing drugs that may cause vitamin D deficiency or interfere with calcium absorption; use of analgesics.", 1],
            ["Surgical Management (Determined by complication severity):", 0, "accent"],
            ["- 1. Fracture repair: The surgical repair of fractures is performed in order to facilitate adequate healing and stability (plates, screws, or rods).", 1],
            ["- 2. Correction of skeletal deformities: Realignment techniques (osteotomy, joint replacement, or spinal fusion).", 1]
        ],
        "images/osteomalacia_legs.png"
    )

    # Slide 23: Osteomalacia: Nursing - Education & Nutrition (verbatim from PDF Page 9)
    add_image_slide(
        prs,
        "Osteomalacia: Nursing",
        "Education & Nutritional Support",
        [
            ["Education:", 0, "primary"],
            ["- 1. Provide patients and their families with thorough information regarding osteomalacia, its causes, and the importance of vitamin D and calcium in bone health.", 1],
            ["- 2. Stress the importance of getting enough sunlight, eating vitamin D-rich foods (such as fatty fish and fortified dairy products), and eating calcium-rich foods.", 1],
            ["- 3. Educate patients on the vitamin D and calcium supplements they have been prescribed, including dosage, administration, and any side effects.", 1],
            ["- 4. Discuss with your health care professional the significance of drug adherence and regular follow-up.", 1],
            ["Nutritional Support:", 0, "accent"],
            ["- 1. Work with dietitians to create tailored vitamin D- and calcium-rich dietary programs.", 1],
            ["- 2. Ensure that patients have access to suitable dietary resources and that they are educated on the food sources of these nutrients.", 1],
            ["- 3. Monitor patients' nutritional status and, if necessary, recommend them for nutritional counseling or supplements.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 24: Osteomalacia: Nursing - Meds & Pain Management (verbatim from PDF Page 9)
    add_image_slide(
        prs,
        "Osteomalacia: Nursing",
        "Medication & Pain Management",
        [
            ["Medication Management:", 0, "primary"],
            ["- 1. Inform patients about the purpose, dosage, administration. and potential side effects of their vitamin D and calcium supplements.", 1],
            ["- 2. Stress the necessity of taking medications as prescribed and taking supplements on a regular basis.", 1],
            ["- 3. Monitor patients for pharmaceutical adverse effects and intervene as needed, or report concerns to the health care team.", 1],
            ["- 4. Teach patients about medication administration time and any food or drug interactions to be cautious of.", 1],
            ["Pain Management:", 0, "accent"],
            ["- 1. Using suitable pain assessment instruments, assess and monitor patients' pain levels on a regular basis.", 1],
            ["- 2. Administer pain drugs as directed and assess their effectiveness.", 1],
            ["- 3. As needed, use nonpharmacologic pain management approaches such as heat or cold therapy, guided visualization, or relaxation exercises.", 1],
            ["- 4. Adjust pain management strategies in collaboration with the health care team based on unique patient needs.", 1]
        ],
        "images/vit_d_metabolism.png"
    )

    # Slide 25: Osteomalacia: Nursing - Mobility & Counseling (verbatim from PDF Page 9 & 10)
    add_image_slide(
        prs,
        "Osteomalacia: Nursing",
        "Mobility & Emotional Support",
        [
            ["Shift Mobility Interventions:", 0, "primary"],
            ["- 1. Evaluate patients' mobility and offer aid and support as needed.", 1],
            ["- 2. Develop tailored exercise routines with physical therapists to improve muscle strength, balance, and mobility.", 1],
            ["- 3. Educate patients on the benefits of regular physical activity and instruct them on safe exercise routines.", 1],
            ["- 4. Assist patients in getting essential assistance aids, such as canes or walkers, to promote safe mobility.", 1],
            ["Counseling and Emotional Support:", 0, "accent"],
            ["- 1. Patients facing pain, functional restrictions, or emotional anguish because of their condition should get emotional support and empathic care.", 1],
            ["- 2. To address patients' anxieties and fears, encourage open conversation and active listening.", 1, "accent"],
            ["- 3. To provide further emotional support, refer patients to suitable options such as support groups or counseling services.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 26: Paget's Disease: Definition & Etiology (verbatim from PDF Page 10)
    add_image_slide(
        prs,
        "Paget's Disease",
        "Definition & Etiology",
        [
            ["Definition:", 0, "primary"],
            ["- Paget's disease of the bone is a skeletal disorder resulting from excessive osteoclast activity, affecting the long bones, pelvis, lumbar vertebrae, and the skull predominantly.", 1, "accent"],
            ["Etiology & Epidemiology:", 0, "primary"],
            ["- 1. The cause of this disease is unknown, although there is evidence of familial tendency (25% to 40% have at least one affected relative).", 1],
            ["- 2. More common in men than in women.", 1],
            ["- 3. Rare before age 40 and increases as age does-12% after age 80.", 1, "accent"],
            ["- 4. May be caused by infection from blood-borne viruses. After acute viremia osteoclasts become chronically infected, stimulating octeoclastic proliferation.", 1]
        ],
        "images/pagets_bone.png"
    )

    # Slide 27: Paget's Disease: Pathophysiology (verbatim from PDF Page 10)
    add_image_slide(
        prs,
        "Paget's Disease",
        "Pathophysiology",
        [
            ["Pathophysiology of Paget's Disease:", 0, "primary"],
            ["- Paget's disease of the bone is a skeletal disorder resulting from excessive osteoclast activity.", 1],
            ["- Infection from blood-borne viruses may cause this disorder.", 1],
            ["- After acute viremia, osteoclasts become chronically infected.", 1, "accent"],
            ["- This chronic infection stimulates osteoclastic proliferation, affecting the long bones, pelvis, lumbar vertebrae, and the skull predominantly.", 1]
        ],
        "images/bone_remodeling.png"
    )

    # Slide 28: Paget's Disease: Clinical Manifestations (verbatim from PDF Page 10)
    add_image_slide(
        prs,
        "Paget's Disease",
        "Clinical Manifestations",
        [
            ["1. Generally asymptomatic.", 0],
            ["2. Most common symptoms are pain and predisposition to fracture.", 0],
            ["3. Pagetic lesions can lead to OA, joint destruction, spinal deformity.", 0, "accent"],
            ["4. Decrease in hearing, tinnitus, and vertigo as a tent of abnormality.", 0],
            ["5. Waddling gait due to abnormality of pelvis.", 0],
            ["6. Radiculopathy and nerve palsies due to effects from de vertebral column.", 0, "accent"],
            ["7. Rarely, heart failure and other cardiovascular effects increased blood supply over abnormal bone.", 0],
            ["8. Malignant bone tumors occur in 5% to 10%", 0]
        ],
        "images/pagets_skull.png"
    )

    # Slide 29: Paget's Disease: Diagnostic Evaluation (verbatim from PDF Page 10)
    add_image_slide(
        prs,
        "Paget's Disease",
        "Diagnostic Evaluation",
        [
            ["1. Elevated serum ALP and urine hydroxypro-line.", 0, "accent"],
            ["2. Serum calcium, phosphorus, and albumin levels usually normal.", 0],
            ["3. Generally confirmed with radiologic examinations showing characteristic abnormalities.", 0],
            ["4. Bone scans can evaluate rapid bone turnover.", 0],
            ["5. Bone biopsy to differentiate from other conditions.", 0]
        ],
        "images/pagets_skull.png"
    )

    # Slide 30: Paget's Disease: Medical & Surgical Management (verbatim from PDF Page 10 & 11)
    add_image_slide(
        prs,
        "Paget's Disease",
        "Medical & Surgical Management",
        [
            ["Medical Management:", 0, "primary"],
            ["- 1. No treatment for asymptomatic Paget's.", 1],
            ["- 2. Pain management-NSAIDs, aspirin.", 1],
            ["- 3. Medications-calcitonin is the main medication used to suppress bone turnover, reduce pain, and prevent progression.", 1, "accent"],
            ["- 4. Other medications used to block bone resorption: bisphosphonates etidronate disodium, alendronate, pamidronate. Risedronate, an antineoplastic agent, plicamycin.", 1],
            ["Surgical Management:", 0, "accent"],
            ["- 5. Tibial osteotomy done to realign knees and relieve pain.", 1]
        ],
        "images/pagets_bone.png"
    )

    # Slide 31: Paget's Disease: Nursing Assessment (verbatim from PDF Page 11)
    add_image_slide(
        prs,
        "Paget's Disease: Nursing",
        "Nursing Assessment",
        [
            ["Nursing Assessment Guidelines:", 0, "primary"],
            ["- 1. Assess pain and functional ability.", 1],
            ["- 2. Observe for bowing (legs) or complaint that hats feel tight.", 1],
            ["- 3. Assess for cardiovascular complications.", 1],
            ["- 4. Assess for auditory symptoms: tinnitus, vertigo, and hearing loss.", 1, "accent"]
        ],
        "images/pagets_skull.png"
    )

    # Slide 32: Paget's Disease: Education & Health Maintenance (verbatim from PDF Page 11)
    add_image_slide(
        prs,
        "Paget's Disease: Nursing",
        "Education & Health Maintenance",
        [
            ["Patient Education and Health Maintenance:", 0, "primary"],
            ["- 1. Teach safety measures in the home-removal of loose rugs and obstacles to prevent falls, good lighting.", 1],
            ["- 2. Provide education about the disease process and medication treatment.", 1],
            ["- 3. Make sure that patient knows how to use mobility aids.", 1],
            ["- 4. Initiate home care referral, as indicated.", 1],
            ["- 5. Provide information about The Paget Foundation (www. paget.org).", 1],
            ["- 6. Encourage follow-up for periodic hearing tests and bloodwork.", 1, "accent"]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 33: Rickets: Introduction & Etiology (verbatim from PDF Page 11)
    add_image_slide(
        prs,
        "Rickets",
        "Introduction & Etiology",
        [
            ["Definition:", 0, "primary"],
            ["- Rickets, a disease of growing bone, occurs in children only before fusion of the epiphyses, and is due to unmineralized matrix at the growth plates.", 1, "accent"],
            ["Etiology:", 0, "primary"],
            ["- Vitamin D disorders: nutritional, congenital, secondary, chronic renal failure, malabsorption.", 1],
            ["- Calcium deficiency: low intake, dietary inhibitors of calcium absorption, malabsorption, premature infants (rickets of prematurity).", 1],
            ["- Phosphorous deficiency: low intake, premature infants (rickets of prematurity), aluminum-containing antacids.", 1],
            ["- Renal losses: X-linked hypophosphatemic rickets, RTA.", 1]
        ],
        "images/rickets_skeletal.png"
    )

    # Slide 34: Rickets: Pathophysiology & Risk Factors (verbatim from PDF Page 11 & 12)
    add_image_slide(
        prs,
        "Rickets",
        "Pathophysiology & Risk Factors",
        [
            ["Pathophysiology:", 0, "primary"],
            ["- 1. Bone consists of a protein matrix called osteoid and a mineral phase, principally composed of calcium and phosphate.", 1],
            ["- 2. Rickets is a disease of growing bone that is caused by unmineralized matrix at the growth plates and occurs in children only before fusion of the epiphyses.", 1],
            ["- 3. Because growth plate cartilage and osteoid continue to expand but mineralization is inadequate, the growth plate thickens.", 1],
            ["- 4. There is also an increase in the circumference of the growth plate and the metaphysis, increasing bone width at the location of the growth plates and causing classic clinical manifestations, such as widening of the wrists and ankles.", 1],
            ["- 5. There is a general softening of the bones that causes them to bend easily when subject to forces such as weight bearing or muscle pull. This softening leads to a variety of bone deformities.", 1, "accent"],
            ["Nutritional Vitamin D Deficiency Risk Factors:", 0, "accent"],
            ["- Living in northern latitudes; dark-skinned people; decreased exposure to sunlight, e.g. in some Asian children living in the UK.", 1],
            ["- Maternal vitamin D deficiency; diets low in calcium, phosphorus and vitamin D, e.g. exclusive breast-feeding into late infancy or toddlers on unsupervised 'dairy-free' diets.", 1],
            ["- Macrobiotic, strict vegan diets; high phytic acid diet, e.g. chapattis.", 1],
            ["- Prolonged parenteral nutrition in infancy with an inadequate supply of parenteral calcium and phosphate.", 1]
        ],
        "images/vit_d_metabolism.png"
    )

    # Slide 35: Rickets: Clinical Signs (General, Head & Chest) (verbatim from PDF Page 12)
    add_image_slide(
        prs,
        "Rickets",
        "Clinical Signs: General, Head & Chest",
        [
            ["General Signs:", 0, "primary"],
            ["- Failure to thrive, listlessness, protruding abdomen, muscle weakness (especially proximal), fractures.", 1],
            ["Head Signs:", 0, "accent"],
            ["- Craniotabes (softening of the cranial bones and can be detected by applying pressure at the occiput or over the parietal bones. The sensation is similar to the feel of pressing into a Ping-Pong ball and then releasing).", 1],
            ["- Frontal bossing, delayed fontanel closure, delayed dentition; caries, craniosynostosis.", 1],
            ["Chest Signs:", 0, "primary"],
            ["- Rachitic rosary, Harrison groove, respiratory infections and atelectasis.", 1]
        ],
        "images/rickets_skeletal.png"
    )

    # Slide 36: Rickets: Clinical Signs (Back & Extremities) (verbatim from PDF Page 12 & 13)
    add_image_slide(
        prs,
        "Rickets",
        "Clinical Signs: Back & Extremities",
        [
            ["Back Signs:", 0, "primary"],
            ["- Scoliosis, kyphosis, lordosis.", 1],
            ["Extremities Signs:", 0, "accent"],
            ["- Enlargement of wrists and ankles, valgus or varus deformities.", 1],
            ["- Windswept deformity (combination of valgus deformity of 1 leg with varus deformity of the other leg).", 1],
            ["- Anterior bowing of the tibia and femur, coxa vara, leg pain.", 1]
        ],
        "images/rickets_skeletal.png"
    )

    # Slide 37: Rickets: Radiologic & Lab Findings (verbatim from PDF Page 13 & 14)
    add_image_slide(
        prs,
        "Rickets",
        "Radiological & Laboratory Findings",
        [
            ["Radiological findings in rickets (best seen at knee, wrist, ankles):", 0, "primary"],
            ["- Frayed, cupped, splaying of the metaphysis.", 1, "accent"],
            ["- Increased distance between the growing epiphysis & metaphysis.", 1],
            ["- Generalized decrease in bone density; bossing and bowing deformity (due to softening); greenstick fractures; rackety rosary of the ribs.", 1],
            ["Laboratory tests (initial tests):", 0, "accent"],
            ["- Serum calcium, phosphorus, alkaline phosphatase (ALP), parathyroid hormone (PTH), 25-hydroxyvitamin D, 1,25-dihydroxyvitamin D3, creatinine, and electrolytes.", 1]
        ],
        "images/bone_types.png"
    )

    # Slide 38: Rickets: Medical & Surgical Management (verbatim from PDF Page 14 & 15)
    add_image_slide(
        prs,
        "Rickets",
        "Medical & Surgical Management",
        [
            ["Medical Treatment:", 0, "primary"],
            ["- Children with nutritional vitamin D deficiency should receive vitamin D and adequate nutritional intake of calcium and phosphorus.", 1],
            ["- Stoss therapy: 300,000-600,000 IU of vitamin D are administered orally or intramuscularly as 2-4 doses over 1 day (ideal if adherence is questionable).", 1],
            ["- Daily, high-dose vitamin D, with doses ranging from 2,000-5,000 IU/day over 4-6 wk.", 1],
            ["- Maintenance: 400 IU/day (<1 yr old) or 600 IU/day (>1 yr old).", 1],
            ["- Symptomatic hypocalcemia: IV calcium acutely, followed by oral calcium supplements tapered over 2-6 wk.", 1],
            ["Surgical Management (Required in severe deformities or when medical treatment fails):", 0, "accent"],
            ["- 1. Corrective osteotomy: surgical correction of bowed legs or deformities.", 1],
            ["- 2. Guided growth surgery: used in growing children to correct angular deformities.", 1],
            ["- 3. External fixation devices: for severe bone deformities.", 1],
            ["- 4. Fracture management: casting or surgical fixation if fractures occur.", 1]
        ],
        "images/rickets_skeletal.png"
    )

    # Slide 39: Rickets: Nursing & Prevention (verbatim from PDF Page 15)
    add_image_slide(
        prs,
        "Rickets",
        "Nursing Management & Prevention",
        [
            ["Nursing Management:", 0, "primary"],
            ["- Pain Management: Provide comfortable positioning; administer analgesics if prescribed.", 1],
            ["- Health Education: Explain disease process, importance of treatment compliance, follow-up visits and monitoring.", 1],
            ["- Psychological Support: Support child and family emotionally; reduce anxiety related to deformities.", 1],
            ["Prevention:", 0, "accent"],
            ["- Adequate sunlight exposure.", 1],
            ["- Proper nutrition during childhood.", 1],
            ["- Vitamin D supplementation in high-risk children.", 1],
            ["- Maternal nutrition during pregnancy.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 40: Comparative Study Table
    add_table_slide(
        prs,
        "Comparative Study",
        "Comparison of Metabolic Bone Disorders",
        ["Disease", "Primary Pathology", "Serum Ca & PO4", "Serum ALP", "Key Treatments"],
        [
            [
                "Osteoporosis", 
                "Bone resorption > bone formation (uncoupled remodeling, low bone mass)", 
                "Normal", 
                "Normal", 
                "Bisphosphonates, Denosumab, Teriparatide, Calcium/Vit D, Exercise"
            ],
            [
                "Osteomalacia", 
                "Inadequate mineralization of organic osteoid matrix (soft bones)", 
                "Low or Low-Normal", 
                "Elevated", 
                "High-dose Vitamin D supplementation, Calcium, treatment of malabsorption"
            ],
            [
                "Paget's Disease", 
                "Chaotic bone remodeling (giant osteoclasts followed by woven bone formation)", 
                "Normal\n(unless immobilized)", 
                "Markedly Elevated", 
                "Bisphosphonates (Zoledronic Acid), Calcitonin, NSAIDs, Decompression surgery"
            ],
            [
                "Rickets", 
                "Inadequate mineralization of growth plate osteoid (children before epiphyseal fusion)", 
                "Low or Low-Normal", 
                "Elevated", 
                "Vitamin D (Stoss or daily), Calcium/Phosphate, guided growth/osteotomy"
            ]
        ],
        "images/metabolic_bone_overview.png"
    )

    # Slide 41: Care Plan 1 - Acute Pain
    add_nursing_care_plan_slide(
        prs,
        "Care Plan 1: Acute Pain",
        [
            "Subjective Data: Patient complains of bone or joint pain. Reports pain during movement.",
            "Objective Data: Facial grimacing, Guarding behavior, Limited mobility, Pain score elevated"
        ],
        "Acute Pain related to bone demineralization and musculoskeletal changes as evidenced by verbal reports of pain and restricted movement.",
        [
            "Patient will report pain ≤3/10 within 48 hours.",
            "Patient will demonstrate improved comfort and mobility."
        ],
        [
            "1. Assess pain intensity, location, and duration regularly.",
            "2. Administer prescribed analgesics.",
            "3. Position patient comfortably and support affected limbs.",
            "4. Encourage relaxation techniques.",
            "5. Apply heat or cold therapy if prescribed."
        ],
        [
            "Patient reports decreased pain. Improved participation in activities."
        ]
    )

    # Slide 42: Care Plan 2 - Impaired Physical Mobility
    add_nursing_care_plan_slide(
        prs,
        "Care Plan 2: Impaired Physical Mobility",
        [
            "Subjective Data: Difficulty walking. Complains of weakness.",
            "Objective Data: Limited range of motion. Requires assistance with movement."
        ],
        "Impaired Physical Mobility related to pain, skeletal deformity, and muscle weakness as evidenced by limited movement and difficulty ambulating.",
        [
            "Patient will demonstrate improved mobility within one week.",
            "Patient will perform activities with minimal assistance."
        ],
        [
            "1. Assess mobility status daily.",
            "2. Encourage active and passive ROM exercises.",
            "3. Assist with ambulation as needed.",
            "4. Collaborate with physiotherapist.",
            "5. Provide assistive devices when indicated."
        ],
        [
            "Patient demonstrates improved mobility. Performs ADLs with less assistance."
        ]
    )

    # Slide 43: Care Plan 3 - Risk for Injury
    add_nursing_care_plan_slide(
        prs,
        "Care Plan 3: Risk for Injury",
        [
            "Risk Factors: Decreased bone density, Muscle weakness, History of falls"
        ],
        "Risk for Injury related to bone fragility and impaired mobility.",
        [
            "Patient will remain free from fractures and injuries during hospitalization."
        ],
        [
            "1. Assess fall risk regularly.",
            "2. Maintain clutter-free environment.",
            "3. Use side rails appropriately.",
            "4. Encourage use of assistive devices.",
            "5. Educate patient on fall-prevention measures."
        ],
        [
            "No falls or injuries reported."
        ]
    )

    # Slide 44: Care Plan 4 - Activity Intolerance
    add_nursing_care_plan_slide(
        prs,
        "Care Plan 4: Activity Intolerance",
        [
            "Subjective Data: Fatigue during activities.",
            "Objective Data: Weakness, Reduced endurance"
        ],
        "Activity Intolerance related to musculoskeletal weakness and pain.",
        [
            "Patient will perform activities without excessive fatigue."
        ],
        [
            "1. Assess response to activity.",
            "2. Schedule rest periods.",
            "3. Gradually increase activity level.",
            "4. Assist with ADLs as necessary.",
            "5. Encourage balanced nutrition."
        ],
        [
            "Patient performs activities with improved tolerance."
        ]
    )

    # Slide 45: Care Plan 5 - Deficient Knowledge
    add_nursing_care_plan_slide(
        prs,
        "Care Plan 5: Deficient Knowledge",
        [
            "Subjective/Objective: Asks questions about disease and treatment."
        ],
        "Deficient Knowledge related to lack of information about metabolic bone disorder and its management.",
        [
            "Patient will verbalize understanding of disease, treatment, diet, and prevention measures."
        ],
        [
            "1. Assess current knowledge level.",
            "2. Explain disease process in simple language.",
            "3. Teach importance of calcium and vitamin D intake.",
            "4. Instruct on medication adherence.",
            "5. Educate regarding exercise and fall prevention."
        ],
        [
            "Patient correctly explains disease management and preventive measures."
        ]
    )

    # Slide 46: Health Education: Nutrition, Sunlight & Exercise (verbatim from PDF Page 19)
    add_image_slide(
        prs,
        "Health Education",
        "Nutrition, Sunlight & Exercise",
        [
            ["1. Nutrition:", 0, "primary"],
            ["- Consume a balanced diet rich in calcium and vitamin D.", 1],
            ["- Include milk, curd, cheese, green leafy vegetables, fish, eggs, and fortified foods.", 1],
            ["- Maintain adequate protein intake for bone health.", 1],
            ["- Avoid excessive salt, caffeine, and carbonated drinks.", 1],
            ["2. Sunlight Exposure:", 0, "accent"],
            ["- Take regular sunlight exposure for 15–30 minutes daily, preferably in the morning.", 1],
            ["- Sunlight helps the body produce vitamin D, which is essential for calcium absorption.", 1],
            ["3. Exercise and Physical Activity:", 0, "primary"],
            ["- Perform regular weight-bearing exercises such as walking, jogging, and stair climbing.", 1],
            ["- Engage in muscle-strengthening exercises as advised by a healthcare provider.", 1],
            ["- Avoid prolonged bed rest and inactivity.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 47: Health Education: Medications, Falls & Lifestyle (verbatim from PDF Page 19)
    add_image_slide(
        prs,
        "Health Education",
        "Medications, Falls & Lifestyle",
        [
            ["4. Medication Adherence:", 0, "primary"],
            ["- Take prescribed calcium, vitamin D supplements, and other medications regularly.", 1],
            ["- Do not stop medications without consulting a doctor.", 1],
            ["- Attend regular follow-up appointments.", 1],
            ["5. Fall Prevention:", 0, "accent"],
            ["- Keep floors dry and free of obstacles.", 1],
            ["- Use proper footwear with non-slip soles.", 1],
            ["- Install handrails in bathrooms and stairways if needed.", 1],
            ["- Use walking aids when recommended.", 1],
            ["6. Lifestyle Modification:", 0, "primary"],
            ["- Avoid smoking and tobacco use.", 1],
            ["- Limit alcohol consumption.", 1],
            ["- Maintain a healthy body weight.", 1]
        ],
        "images/nursing_care_icon.png"
    )

    # Slide 48: Health Education: Monitoring & Family Education (verbatim from PDF Page 19)
    add_image_slide(
        prs,
        "Health Education",
        "Monitoring & Family Education",
        [
            ["7. Monitoring and Follow-up:", 0, "primary"],
            ["- Regularly monitor calcium, phosphate, and vitamin D levels as advised.", 1],
            ["- Undergo bone density testing when recommended.", 1],
            ["- Report symptoms such as bone pain, fractures, muscle weakness, or deformities promptly.", 1],
            ["8. Patient and Family Education:", 0, "accent"],
            ["- Understand the nature of the disease and treatment plan.", 1],
            ["- Recognize signs of complications.", 1],
            ["- Encourage family support for dietary and lifestyle modifications.", 1]
        ],
        "images/metabolic_bone_overview.png"
    )

    # Slide 49: Conclusion (verbatim from PDF Page 20)
    add_image_slide(
        prs,
        "Conclusion",
        "Clinical Summary & Multidisciplinary Role",
        [
            ["Clinical Summary:", 0, "primary"],
            ["- Metabolic bone disorders are a group of conditions that affect bone strength, structure, and mineralization due to abnormalities in calcium, phosphate, vitamin D metabolism, or hormonal regulation. Common disorders include Osteoporosis, Osteomalacia, and Rickets.", 1],
            ["- These conditions can lead to bone pain, deformities, fractures, reduced mobility, and impaired quality of life. Early diagnosis through clinical assessment, laboratory investigations, and imaging studies is essential for effective management.", 1],
            ["Treatment Focus:", 0, "accent"],
            ["- Focuses on correcting underlying causes, restoring mineral balance, preventing complications, and promoting bone health through medications, nutritional support, lifestyle modifications, and patient education.", 1],
            ["Multidisciplinary Role:", 0, "primary"],
            ["- A multidisciplinary approach involving physicians, nurses, dietitians, and physiotherapists plays a vital role in improving patient outcomes and preventing long-term disability.", 1]
        ],
        "images/metabolic_bone_overview.png"
    )

    # Slide 50: Bibliography (verbatim from PDF Page 20)
    add_image_slide(
        prs,
        "References",
        "Bibliography & Key Textbooks",
        [
            ["Clinical Textbooks & References:", 0, "primary"],
            ["- Brunner & Suddarth's Textbook of Medical-Surgical Nursing. Hinkle JL, Cheever KH. Philadelphia: Wolters Kluwer; 2022.", 1],
            ["- Lewis's Medical-Surgical Nursing. Harding MM, Kwong J, Roberts D, Hagler D. Elsevier; 2023.", 1],
            ["- Textbook of Medical-Surgical Nursing. Williams & Wilkins; Latest Edition.", 1],
            ["- Essentials of Medical Surgical Nursing. I. Clement. Latest Edition.", 1],
            ["- Medical-Surgical Nursing. CBS Publishers & Distributors; Latest Edition.", 1]
        ],
        "images/metabolic_bone_intro.png"
    )

    # Save
    output_filename = "Metabolic_Bone_Disorders.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}' in workspace.")

if __name__ == "__main__":
    main()
